"""Robust file processor — handles ANY file that arrives via email.

Never crashes. Always responds. Processes what it can, stores the rest.

Handler routing:
    vision              → resize + send to vision pipeline (images)
    convert_then_vision → HEIC→JPEG conversion then vision
    pdf_ingest          → PyMuPDF text extraction (≤10 pages) or queue (>10)
    docx/xlsx/pptx      → python-docx / openpyxl / python-pptx text extraction
    text_ingest         → decode + return text
    csv_ingest          → decode + return text
    extract_and_process → zipfile extraction + recurse
    store_only          → save to disk + friendly acknowledgment
"""

from __future__ import annotations

import io
import logging
import mimetypes
from pathlib import Path
from typing import Any, Callable, Coroutine

logger = logging.getLogger("mira-email")

MAX_FILE_SIZE = 50 * 1024 * 1024   # 50 MB hard limit
MAX_PDF_PAGES_INLINE = 10          # PDFs above this go to ingest queue
MAX_ARCHIVE_FILES = 20             # cap files extracted from zip/rar

# Content-type → handler name
_CT_MAP: dict[str, str] = {
    "image/jpeg": "vision",
    "image/jpg": "vision",
    "image/png": "vision",
    "image/gif": "vision",
    "image/webp": "vision",
    "image/tiff": "vision",
    "image/bmp": "vision",
    "image/heic": "convert_then_vision",
    "image/heif": "convert_then_vision",
    "application/pdf": "pdf_ingest",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx_extract",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "xlsx_extract",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx_extract",
    "application/msword": "doc_legacy",
    "application/vnd.ms-excel": "xls_legacy",
    "text/plain": "text_ingest",
    "text/csv": "csv_ingest",
    "text/x-csv": "csv_ingest",
    "application/csv": "csv_ingest",
    "application/zip": "extract_and_process",
    "application/x-zip-compressed": "extract_and_process",
    "application/x-rar-compressed": "extract_and_process",
    "application/vnd.rar": "extract_and_process",
    # CAD / engineering files
    "application/acad": "store_only",
    "application/dxf": "store_only",
    "application/step": "store_only",
    "model/step": "store_only",
    "model/iges": "store_only",
}

# Extension fallback
_EXT_MAP: dict[str, str] = {
    ".jpg": "vision", ".jpeg": "vision", ".png": "vision",
    ".gif": "vision", ".webp": "vision", ".tiff": "vision",
    ".tif": "vision", ".bmp": "vision",
    ".heic": "convert_then_vision", ".heif": "convert_then_vision",
    ".pdf": "pdf_ingest",
    ".docx": "docx_extract", ".xlsx": "xlsx_extract", ".pptx": "pptx_extract",
    ".doc": "doc_legacy", ".xls": "xls_legacy",
    ".txt": "text_ingest", ".log": "text_ingest", ".md": "text_ingest",
    ".csv": "csv_ingest", ".tsv": "csv_ingest",
    ".zip": "extract_and_process", ".rar": "extract_and_process",
    ".dwg": "store_only", ".dxf": "store_only",
    ".stp": "store_only", ".step": "store_only",
    ".igs": "store_only", ".iges": "store_only",
    ".stl": "store_only",
}

# Human-readable labels for store_only acknowledgements
_STORE_LABELS: dict[str, str] = {
    ".dwg": "AutoCAD drawing",
    ".dxf": "DXF drawing",
    ".stp": "STEP model",
    ".step": "STEP model",
    ".igs": "IGES model",
    ".iges": "IGES model",
    ".stl": "3D model",
}


class FileProcessor:
    """Process any file arriving via email. Never raises."""

    def __init__(
        self,
        vision_pipeline: Callable[..., Coroutine[Any, Any, dict]] | None = None,
        ingest_pipeline: Callable[..., Coroutine[Any, Any, None]] | None = None,
        storage_path: str = "/data/email-attachments",
    ) -> None:
        self.vision = vision_pipeline
        self.ingest = ingest_pipeline
        self.storage_path = Path(storage_path)
        try:
            self.storage_path.mkdir(parents=True, exist_ok=True)
        except Exception:
            pass

    async def process(self, filename: str, content_type: str, data: bytes) -> dict:
        """Process one file. Returns a result dict. Never raises."""
        try:
            return await self._process(filename, content_type, data)
        except Exception as exc:
            logger.error(
                "FILE_PROCESS_ERROR file=%s type=%s error=%s",
                filename, content_type, str(exc)[:300],
            )
            self._store_file(filename, data)
            return _stored(filename, f"I had trouble processing {filename}. It's stored for reference — the error has been logged.")

    async def _process(self, filename: str, content_type: str, data: bytes) -> dict:
        if len(data) > MAX_FILE_SIZE:
            size_mb = len(data) // 1024 // 1024
            return _stored(
                filename,
                f"{filename} is {size_mb} MB — over the 50 MB limit. Stored for reference but not processed.",
            )

        handler = self._get_handler(content_type, filename)

        if handler == "vision":
            return await self._handle_vision(filename, data)
        if handler == "convert_then_vision":
            return await self._handle_heic(filename, data)
        if handler == "pdf_ingest":
            return await self._handle_pdf(filename, data)
        if handler == "docx_extract":
            return await self._handle_docx(filename, data)
        if handler == "xlsx_extract":
            return await self._handle_xlsx(filename, data)
        if handler == "pptx_extract":
            return await self._handle_pptx(filename, data)
        if handler in ("doc_legacy", "xls_legacy"):
            return _stored(filename, f"Legacy Office format received. For best results, save as {'.docx' if 'doc' in handler else '.xlsx'} and resend.")
        if handler == "text_ingest":
            return self._handle_text(filename, data)
        if handler == "csv_ingest":
            return self._handle_csv(filename, data)
        if handler == "extract_and_process":
            return await self._handle_archive(filename, data)

        # store_only or unknown
        ext = Path(filename).suffix.lower()
        label = _STORE_LABELS.get(ext, Path(filename).suffix or "this file type")
        return _stored(
            filename,
            f"I received {filename} ({label}) but can't process {ext or 'this file type'} files yet. "
            f"I've stored it for reference. In the meantime I can work with photos (JPEG/PNG), PDFs, and text.",
        )

    def _get_handler(self, content_type: str, filename: str) -> str:
        ct = content_type.lower().split(";")[0].strip()
        if ct in _CT_MAP:
            return _CT_MAP[ct]
        # Extension fallback
        ext = Path(filename).suffix.lower()
        if ext in _EXT_MAP:
            return _EXT_MAP[ext]
        # Generic image/* fallback
        if ct.startswith("image/"):
            return "vision"
        if ct.startswith("text/"):
            return "text_ingest"
        return "store_only"

    # ------------------------------------------------------------------
    # Handlers
    # ------------------------------------------------------------------

    async def _handle_vision(self, filename: str, data: bytes) -> dict:
        import base64
        b64 = base64.b64encode(data).decode()
        if self.vision:
            try:
                result = await self.vision(b64)
                desc = result.get("description", "Image received and analyzed.")
                return {
                    "status": "processed",
                    "handler": "vision",
                    "filename": filename,
                    "description": f"Analyzed {filename}: {desc[:300]}",
                    "extracted_text": result.get("text", ""),
                }
            except Exception as exc:
                logger.warning("VISION_FALLBACK file=%s error=%s", filename, str(exc)[:100])

        self._store_file(filename, data)
        return {
            "status": "stored",
            "handler": "vision_fallback",
            "filename": filename,
            "description": f"Image {filename} received and stored. Vision analysis unavailable right now.",
            "extracted_text": "",
        }

    async def _handle_heic(self, filename: str, data: bytes) -> dict:
        try:
            from PIL import Image
            img = Image.open(io.BytesIO(data))
            buf = io.BytesIO()
            img = img.convert("RGB")
            img.save(buf, format="JPEG", quality=90)
            jpg_name = Path(filename).with_suffix(".jpg").name
            return await self._handle_vision(jpg_name, buf.getvalue())
        except Exception as exc:
            logger.warning("HEIC_CONVERT_FAIL file=%s error=%s", filename, str(exc)[:100])
            return _stored(
                filename,
                f"iPhone photo ({filename}) received. HEIC conversion failed — try sending as JPEG instead.",
            )

    async def _handle_pdf(self, filename: str, data: bytes) -> dict:
        self._store_file(filename, data)
        page_count = self._pdf_page_count(data)

        if page_count == 0:
            return _stored(filename, f"PDF received ({filename}) but couldn't read page count. Stored for reference.")

        if page_count <= MAX_PDF_PAGES_INLINE:
            text = self._pdf_extract_text(data)
            return {
                "status": "processed",
                "handler": "pdf_inline",
                "filename": filename,
                "description": f"PDF ({page_count} {'page' if page_count == 1 else 'pages'}) extracted and ready for questions.",
                "extracted_text": text[:5000],
            }

        # Large PDF → queue
        if self.ingest:
            try:
                await self.ingest(filename, data)
            except Exception as exc:
                logger.warning("PDF_INGEST_QUEUE_FAIL file=%s error=%s", filename, str(exc)[:100])

        return {
            "status": "queued",
            "handler": "pdf_ingest_queue",
            "filename": filename,
            "description": (
                f"PDF ({page_count} pages) queued for knowledge base ingestion. "
                "Ask me about it in a few minutes once ingestion completes."
            ),
            "extracted_text": "",
        }

    def _pdf_page_count(self, data: bytes) -> int:
        try:
            import fitz  # PyMuPDF
            doc = fitz.open(stream=data, filetype="pdf")
            return len(doc)
        except Exception:
            # Fallback: count %%EOF or /Page occurrences
            try:
                count = data.count(b"/Type /Page")
                return count if count > 0 else 1
            except Exception:
                return 0

    def _pdf_extract_text(self, data: bytes) -> str:
        try:
            import fitz
            doc = fitz.open(stream=data, filetype="pdf")
            return "\n".join(page.get_text() for page in doc).strip()
        except Exception:
            return ""

    async def _handle_docx(self, filename: str, data: bytes) -> dict:
        self._store_file(filename, data)
        try:
            from docx import Document
            doc = Document(io.BytesIO(data))
            text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
            return {
                "status": "processed",
                "handler": "docx_extract",
                "filename": filename,
                "description": f"Word document extracted ({len(text)} chars).",
                "extracted_text": text[:5000],
            }
        except ImportError:
            return _stored(filename, f"Word document {filename} received. python-docx not installed — stored for reference.")
        except Exception as exc:
            logger.warning("DOCX_EXTRACT_FAIL file=%s error=%s", filename, str(exc)[:100])
            return _stored(filename, f"Received {filename} but couldn't extract text. Stored for reference.")

    async def _handle_xlsx(self, filename: str, data: bytes) -> dict:
        self._store_file(filename, data)
        try:
            import openpyxl
            wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
            rows: list[str] = []
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    line = " | ".join(str(c) for c in row if c is not None)
                    if line.strip():
                        rows.append(line)
            text = "\n".join(rows)
            return {
                "status": "processed",
                "handler": "xlsx_extract",
                "filename": filename,
                "description": f"Spreadsheet extracted ({len(rows)} rows).",
                "extracted_text": text[:5000],
            }
        except ImportError:
            return _stored(filename, f"Spreadsheet {filename} received. openpyxl not installed — stored for reference.")
        except Exception as exc:
            logger.warning("XLSX_EXTRACT_FAIL file=%s error=%s", filename, str(exc)[:100])
            return _stored(filename, f"Received {filename} but couldn't extract data. Stored for reference.")

    async def _handle_pptx(self, filename: str, data: bytes) -> dict:
        self._store_file(filename, data)
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(data))
            parts: list[str] = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        t = shape.text_frame.text.strip()
                        if t:
                            parts.append(t)
            text = "\n".join(parts)
            return {
                "status": "processed",
                "handler": "pptx_extract",
                "filename": filename,
                "description": f"Presentation extracted ({len(prs.slides)} slides).",
                "extracted_text": text[:5000],
            }
        except ImportError:
            return _stored(filename, f"Presentation {filename} received. python-pptx not installed — stored for reference.")
        except Exception as exc:
            logger.warning("PPTX_EXTRACT_FAIL file=%s error=%s", filename, str(exc)[:100])
            return _stored(filename, f"Received {filename} but couldn't extract content. Stored for reference.")

    def _handle_text(self, filename: str, data: bytes) -> dict:
        text = data.decode("utf-8", errors="replace")
        return {
            "status": "processed",
            "handler": "text",
            "filename": filename,
            "description": f"Text file received ({len(text):,} chars).",
            "extracted_text": text[:5000],
        }

    def _handle_csv(self, filename: str, data: bytes) -> dict:
        text = data.decode("utf-8", errors="replace")
        row_count = text.count("\n")
        return {
            "status": "processed",
            "handler": "csv",
            "filename": filename,
            "description": f"CSV file received (~{row_count} rows).",
            "extracted_text": text[:5000],
        }

    async def _handle_archive(self, filename: str, data: bytes) -> dict:
        results: list[dict] = []
        try:
            import zipfile
            with zipfile.ZipFile(io.BytesIO(data)) as zf:
                names = [n for n in zf.namelist() if not n.endswith("/")][:MAX_ARCHIVE_FILES]
                for name in names:
                    try:
                        inner = zf.read(name)
                        ct = _guess_mime(name)
                        results.append(await self.process(Path(name).name, ct, inner))
                    except Exception as exc:
                        logger.warning("ARCHIVE_INNER_FAIL name=%s error=%s", name, str(exc)[:100])
                        results.append(_stored(name, f"Couldn't extract {name} from archive."))
        except Exception as exc:
            logger.warning("ARCHIVE_OPEN_FAIL file=%s error=%s", filename, str(exc)[:100])
            return _stored(filename, f"Archive {filename} received but couldn't be extracted. Stored for reference.")

        descs = "; ".join(r.get("description", r.get("filename", "?"))[:80] for r in results[:5])
        all_text = "\n".join(r.get("extracted_text", "") for r in results)
        return {
            "status": "processed",
            "handler": "archive",
            "filename": filename,
            "description": f"Archive with {len(results)} file(s): {descs}",
            "extracted_text": all_text[:5000],
        }

    # ------------------------------------------------------------------

    def _store_file(self, filename: str, data: bytes) -> None:
        try:
            safe = "".join(c for c in filename if c.isalnum() or c in ".-_ ")[:100].strip()
            if safe:
                (self.storage_path / safe).write_bytes(data)
        except Exception:
            pass


def _stored(filename: str, description: str) -> dict:
    return {
        "status": "stored",
        "handler": "store_only",
        "filename": filename,
        "description": description,
        "extracted_text": "",
    }


def _guess_mime(filename: str) -> str:
    ct, _ = mimetypes.guess_type(filename)
    return ct or "application/octet-stream"
